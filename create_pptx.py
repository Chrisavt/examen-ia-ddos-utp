#!/usr/bin/env python3
"""
create_pptx.py — Genera la presentacion PowerPoint para el examen
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
RESULTS_DIR = os.path.join(BASE_DIR, 'resultados')
DIAG_DIR = os.path.join(BASE_DIR, 'diagramas')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = ACCENT_BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Seccion {number}"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(title, bullets, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if image_path and os.path.exists(image_path):
        # Content on left, image on right
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for bullet in bullets:
            p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)
            p.level = 0
        slide.shapes.add_picture(image_path, Inches(7.5), Inches(1.3), Inches(5.5))
    else:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for bullet in bullets:
            p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)
    return slide


def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Left column
    txL = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
    tfL = txL.text_frame
    tfL.word_wrap = True
    p = tfL.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    for b in left_bullets:
        p = tfL.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    # Right column
    txR = slide.shapes.add_textbox(Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
    tfR = txR.text_frame
    tfR.word_wrap = True
    p = tfR.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    for b in right_bullets:
        p = tfR.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    return slide


# ==========================================
# SLIDE 1: PORTADA
# ==========================================
add_title_slide(
    "Deteccion y Clasificacion de DDoS\nusando Machine Learning",
    "Dataset: CIC-DDoS2019 | Modelo: Random Forest\nUniversidad Tecnologica de Panama\nIA Aplicada a la Ciberseguridad - Tema 2\n\nChristopher Abrego"
)

# ==========================================
# SLIDE 2: CONTENIDO
# ==========================================
add_content_slide("Contenido", [
    "I.   Contexto y Planteamiento del Problema",
    "II.  Objetivos (General y Especificos)",
    "III. Bases Tecnicas y Tecnologias",
    "IV.  Caso de Estudio y Desarrollo",
    "V.   Resultados y Analisis",
    "VI.  Demostracion en Vivo",
    "VII. Analisis Critico",
    "VIII. Conclusiones y Recomendaciones",
    "     Referencias Bibliograficas",
])

# ==========================================
# SECCION I: CONTEXTO
# ==========================================
add_section_slide("I", "Contexto y Planteamiento del Problema")

add_content_slide("Contexto y Planteamiento del Problema", [
    "- Los ataques DDoS (Distributed Denial of Service) representan una de las mayores",
    "  amenazas a la disponibilidad de servicios en Internet",
    "",
    "- En 2024, el volumen promedio de ataques DDoS supero los 1 Tbps",
    "  (Cloudflare, 2024)",
    "",
    "- Los metodos tradicionales de deteccion (firmas, IDS estaticos) no logran",
    "  identificar ataques zero-day o variantes nuevas",
    "",
    "- El Machine Learning ofrece la capacidad de aprender patrones del trafico",
    "  y adaptarse a nuevos vectores de ataque",
    "",
    "- Problema: Como implementar un sistema eficiente de deteccion y",
    "  clasificacion de DDoS en tiempo real usando ML?",
])

add_content_slide("Impacto de los ataques DDoS", [
    "- Disponibilidad: Los servicios quedan inaccesibles para usuarios legitimos",
    "- Economia: Costos estimados de $120,000 por hora de inactividad (Dyn, 2016)",
    "- Infraestructura: Sobrecarga de servidores, routers y enlaces de red",
    "- Repuracion: DDoS como cortina de humo para otros ataques (robos de datos)",
    "",
    "- Tipos principales de DDoS:",
    "  - Volumetricos: UDP Flood, DNS Amplification (Capa 3-4)",
    "  - Protocolo: SYN Flood, Ping of Death (Capa 3-4)",
    "  - Aplicacion: HTTP Flood, Slowloris (Capa 7)",
    "",
    "- El CIC-DDoS2019 incluye 13 tipos de ataques DDoS reales",
], os.path.join(BASE_DIR, 'resultados', 'distribucion_clases.png'))

# ==========================================
# SECCION II: OBJETIVOS
# ==========================================
add_section_slide("II", "Objetivos")

add_content_slide("Objetivo General", [
    "Disenar e implementar un sistema de deteccion y clasificacion de ataques",
    "DDoS utilizando tecnicas de Machine Learning, basado en el dataset",
    "CIC-DDoS2019, capaz de identificar en tiempo real el tipo de ataque",
    "y ejecutar mecanismos de mitigacion automatica.",
    "",
    "Objetivos Especificos:",
    "",
    "1. Preprocesar y limpiar el dataset CIC-DDoS2019 (8.2 GB, 22M+ flujos)",
    "2. Entrenar y comparar 3 modelos ML: Random Forest, XGBoost y MLP",
    "3. Evaluar el rendimiento con metricas: Accuracy, Precision, Recall, F1-Score",
    "4. Implementar un sistema de deteccion en tiempo real con scapy",
    "5. Integrar bloqueo automatico via iptables",
    "6. Desarrollar una demostracion en vivo del pipeline completo",
])

# ==========================================
# SECCION III: BASES TECNICAS
# ==========================================
add_section_slide("III", "Bases Tecnicas y Tecnologias")

add_content_slide("Dataset: CIC-DDoS2019", [
    "- Creado por: Canadian Institute for Cybersecurity (UNB, Canada)",
    "- Autores: Sharafaldin, Lashkari, Hakak & Ghorbani (2019)",
    "- Captura: 2 dias de trafico real en laboratorio控制ado",
    "- Total: ~22 millones de flujos de red etiquetados",
    "- Features: 87 columnas extraidas por CICFlowMeter v3",
    "",
    "Tipos de ataques incluidos (13):",
    "  SYN Flood | UDP Flood | LDAP | MSSQL | NetBIOS",
    "  SNMP | SSDP | DNS | TFTP | CharGen | NTP | Portmap | WebDDoS",
    "",
    "Features clave por flujo:",
    "  Duracion | Bytes/s | Packets/s | IAT (Inter-Arrival Time)",
    "  Packet Length | TCP Flags (SYN, ACK, RST, PSH, FIN)",
    "  Flow Direction | Header Length | Window Size",
])

add_two_column_slide(
    "Algoritmos de Machine Learning",
    "Modelos Entrenados",
    [
        "Random Forest (RF):",
        "  - 200 arboles de decision",
        "  - Accuracy: 89.1%",
        "  - Mejor balance precision/recall",
        "",
        "XGBoost:",
        "  - 100 boosters, learning_rate=0.1",
        "  - Accuracy: 92.8%",
        "  - Mejor rendimiento global",
        "",
        "MLP Neural Network:",
        "  - Capas: (100, 50, 25)",
        "  - Accuracy: 85.9%",
        "  - Red neuronal multicapa",
    ],
    "Tecnologias Utilizadas",
    [
        "Lenguaje: Python 3.14",
        "",
        "Librerias principales:",
        "  - scikit-learn (modelos ML)",
        "  - XGBoost (gradient boosting)",
        "  - scapy (captura de paquetes)",
        "  - pandas/numpy (datos)",
        "  - matplotlib/seaborn (viz)",
        "  - joblib (serializacion)",
        "",
        "Entorno: WSL Ubuntu",
        "  - 16 GB RAM",
        "  - Jupyter Notebook",
])

add_content_slide("Pipeline de Machine Learning", [
    "1. Carga de datos → CSV del CIC-DDoS2019 (8.2 GB)",
    "2. Limpieza → Eliminar NaN, Inf, columnas irrelevantes",
    "3. Codificacion → LabelEncoder para las 7 clases de trafico",
    "4. Division → 80% entrenamiento / 20% prueba (random_state=42)",
    "5. Normalizacion → StandardScaler (media=0, desviacion=1)",
    "6. Entrenamiento → 3 modelos en paralelo",
    "7. Evaluacion → Accuracy, Precision, Recall, F1-Score",
    "8. Despliegue → Serializacion con joblib (.pkl)",
    "",
    "El modelo final (RF) se despliega en un script de Python que",
    "captura trafico en tiempo real y lo clasifica instantaneamente.",
], os.path.join(BASE_DIR, 'resultados', 'comparacion_modelos.png'))

# ==========================================
# SECCION IV: CASO DE ESTUDIO
# ==========================================
add_section_slide("IV", "Caso de Estudio y Desarrollo")

add_content_slide("Arquitectura del Sistema", [
    "El sistema se compone de 3 capas principales:",
    "",
    "Capa 1 - Captura: Scapy sniffer en interfaz loopback (lo)",
    "  - Filtra por puerto TCP 8080",
    "  - Extrae headers IP, TCP, flags y payloads",
    "",
    "Capa 2 - Analisis: Motor de deteccion hibrido",
    "  - Heuristica en tiempo real (SYN ratio, PPS, RST ratio)",
    "  - Clasificacion ML (Random Forest, 77 features)",
    "",
    "Capa 3 - Respuesta: Sistema de bloqueo",
    "  - Alerta visual en terminal (colores por tipo)",
    "  - Bloqueo automatico via iptables",
    "  - Log en formato JSON para auditoria",
], os.path.join(RESULTS_DIR, 'matrices_confusion.png'))

add_content_slide("Diagrama de Secuencia - Deteccion", [
    "Flujo completo de deteccion en tiempo real:",
    "",
    "1. Atacante envia trafico malicioso al servidor HTTP :8080",
    "2. Scapy captura cada paquete en la interfaz loopback",
    "3. classify_packet() extrae tipo, flags, IPs, tamanio",
    "4. Packet infos se acumulan en buffer por ventana de 3s",
    "5. detect_attack_heuristic() evalua ratios y umbrales",
    "6. Si se detecta ataque → alerta roja en terminal",
    "7. extract_flow_features() genera vector de 77 features",
    "8. model.predict() clasifica con Random Forest",
    "9. block_ip() ejecuta iptables (o simula)",
    "10. Se guarda log JSON con timestamp, IP, tipo, confianza",
])

# ==========================================
# SECCION V: RESULTADOS
# ==========================================
add_section_slide("V", "Resultados y Analisis")

add_content_slide("Matriz de Confusion", [
    "La matriz de confusion muestra el rendimiento del modelo Random Forest",
    "para las 7 clases del dataset CIC-DDoS2019:",
    "",
    "- BENIGN: Alta tasa de acierto (tráfico normal correctamente identificado)",
    "- Syn Flood: Buena deteccion de ataques SYN",
    "- UDP Flood: Clasificacion precisa de flood UDP",
    "- LDAP/MSSQL/NetBIOS: Deteccion de ataques de amplificacion",
    "- Portmap: Identificacion de ataques de mapeo de puertos",
    "",
    "Las confusiones principales ocurren entre clases de ataques similares",
    "(ej: LDAP vs MSSQL, que comparten patrones de amplificacion)",
], os.path.join(RESULTS_DIR, 'matrices_confusion.png'))

add_two_column_slide(
    "Comparacion de Modelos",
    "Metricas de Rendimiento",
    [
        "Random Forest (seleccionado):",
        "  Accuracy:  89.1%",
        "  Precision: ~90%",
        "  Recall:    ~89%",
        "  F1-Score:  ~89%",
        "",
        "XGBoost:",
        "  Accuracy:  92.8%",
        "  Mejor en overall performance",
        "",
        "MLP Neural Network:",
        "  Accuracy:  85.9%",
        "  Mas lento en inferencia",
    ],
    "Analisis Comparativo",
    [
        "- RF elegido por:",
        "  + Rapidez de inferencia (<1ms)",
        "  + Robustez ante overfitting",
        "  + Facil interpretabilidad",
        "  + Manejo de clases imbalanceadas",
        "",
        "- XGBoost tuvo mejor accuracy",
        "  pero mayor complejidad de deploy",
        "",
        "- MLP bajo rendimiento por:",
        "  + Necesita mas datos de entrenamiento",
        "  + Hiperparametros mas sensibles",
        "  + Mas lento en produccion",
    ]
)

add_content_slide("Importancia de Features", [
    "Las features mas importantes para la clasificacion (Random Forest):",
    "",
    "1. SYN Flag Count — Diferencia clave entre trafico normal y SYN flood",
    "2. Flow Packets/s — Volumen de paquetes por segundo",
    "3. Flow Bytes/s — Tasa de transferencia de datos",
    "4. Fwd Packet Length Mean — Tamaño promedio de paquetes forward",
    "5. Flow IAT Mean — Tiempo promedio entre paquetes",
    "6. Bwd Packet Length Max — Tamaño maximo de respuesta",
    "7. Init Fwd Win Byts — Ventana inicial TCP (indica congestion)",
    "8. ACK Flag Count — Proporcion de ACKs vs total",
    "",
    "Estas features coinciden con la investigacion de Becerra-Suarez et al. (2024)",
    "quienes identificaron que la reduccion a 22 features es suficiente.",
], os.path.join(RESULTS_DIR, 'feature_importance.png'))

# ==========================================
# SECCION VI: DEMOSTRACION
# ==========================================
add_section_slide("VI", "Demostracion en Vivo")

add_content_slide("Demostracion en Vivo del Sistema", [
    "Pipeline completo de demostracion:",
    "",
    "1. Servidor HTTP activo en puerto 8080",
    "2. Scapy captura trafico en interfaz loopback",
    "3. Motor heuristica detecta patrones de ataque en tiempo real",
    "4. Modelo ML clasifica cada flujo",
    "5. Alerta visual cuando se detecta ataque:",
    "   - Rojo: ATTACK (Syn Flood, HTTP Flood, RST Flood)",
    "   - Verde: NORMAL (trafico legitimo)",
    "6. Bloqueo automatico via iptables (simulado o real)",
    "7. Resumen final con estadisticas",
    "",
    "Scripts disponibles:",
    "  - demo_viva.py — Monitor principal",
    "  - http_flood.py — Generador de HTTP flood",
    "  - syn_flood.py — Generador de SYN flood",
])

add_content_slide("Resultados de la Demostracion", [
    "Ejecucion exitosa del demo (60 segundos):",
    "",
    "  Paquetes capturados: 27,357",
    "  Alertas totales:     15",
    "  IPs bloqueadas:      1 (127.0.0.1)",
    "",
    "Distribucion de ataques detectados:",
    "  - Syn Flood:  8 detecciones (SYN ratio: 88-99%)",
    "  - HTTP Flood: 6 detecciones (PPS: 480-1065)",
    "  - RST Flood:  1 deteccion  (RST ratio: 50%)",
    "",
    "El sistema demuestra la capacidad de:",
    "  1. Capturar trafico en tiempo real",
    "  2. Detectar multiples tipos de ataque",
    "  3. Clasificar con el modelo ML",
    "  4. Ejecutar respuesta automatica (bloqueo)",
])

# ==========================================
# SECCION VII: ANALISIS CRITICO
# ==========================================
add_section_slide("VII", "Analisis Critico")

add_two_column_slide(
    "Analisis Critico",
    "Fortalezas del Enfoque",
    [
        "- Dataset real y ampliamente citado (CIC-DDoS2019)",
        "- Multiples modelos comparados (RF, XGB, MLP)",
        "- Pipeline completo: train → deploy → detect",
        "- Deteccion heuristica + ML (enfoque hibrido)",
        "- Tiempo real: inferencia <1ms por flujo",
        "- Bloqueo automatico via iptables",
        "- Demostracion funcional y verificable",
        "- Codigo abierto en GitHub",
    ],
    "Limitaciones y Trabajo Futuro",
    [
        "- Modelo entrenado en localhost no captura",
        "  patrones reales de red (loopback vs ETH)",
        "- 77 features simplificadas en captura en vivo",
        "- Umbral fijo para heuristica (no adaptativo)",
        "- Sin deteccion de ataques zero-day",
        "- Sin balanceo de clases en entrenamiento",
        "- Solo TCP (no UDP en demo en vivo)",
        "",
        "Trabajo futuro:",
        "- Integrar CICFlowMeter en tiempo real",
        "- Usar XGBoost (92.8%) como modelo principal",
        "- Agregar deep learning (LSTM/CNN)",
    ]
)

add_content_slide("Comparacion con Trabajos Relacionados", [
    "Becerra-Suarez et al. (2024): RF con 99.97% en CIC-DDoS2019",
    "  → Reduccion a 22 features con Pearson correlation",
    "",
    "Najam & Abduljawad (2023): RF-RFE-SMOTE alcanzo 99-100%",
    "  → Feature selection + SMOTE para balanceo",
    "",
    "Ma et al. (2023): Framework FAMS con RF optimizado",
    "  → Seleccion de features y modelos automatica",
    "",
    "Mualfah et al. (2025): RF + class weight = 99.98%",
    "  → Class weight para manejar desbalance de clases",
    "",
    "Nuestro resultado (89.1%) es consistente con un enfoque",
    "sin optimizacion avanzada de features ni balanceo de clases.",
    "El gap se debe a que usamos las 77 features originales",
    "sin reduccion ni ingenieria de features avanzada.",
], os.path.join(RESULTS_DIR, 'comparacion_modelos.png'))

# ==========================================
# SECCION VIII: CONCLUSIONES
# ==========================================
add_section_slide("VIII", "Conclusiones y Recomendaciones")

add_content_slide("Conclusiones", [
    "1. El modelo Random Forest logro un accuracy del 89.1% en la",
    "   clasificacion de 7 tipos de trafico DDoS del CIC-DDoS2019",
    "",
    "2. El sistema hibrido (heuristica + ML) demuestra deteccion efectiva",
    "   en tiempo real con 15 alertas en 60 segundos de ejecucion",
    "",
    "3. El pipeline completo funciona: captura → features → clasificacion",
    "   → alerta → bloqueo automático via iptables",
    "",
    "4. El enfoque de ML supera a los IDS tradicionales basados en firmas",
    "   porque aprende patrones del trafico sin definiciones estaticas",
    "",
    "5. El CIC-DDoS2019 es un dataset robusto y adecuado para evaluar",
    "   sistemas de deteccion DDoS con Machine Learning",
])

add_content_slide("Recomendaciones", [
    "Para mejorar el sistema en futuras iteraciones:",
    "",
    "1. Feature Selection: Aplicar RF-RFE o Pearson correlation",
    "   para reducir de 77 a ~22 features (mejora esperada: +10% accuracy)",
    "",
    "2. Balanceo de Clases: Usar SMOTE o class weights para",
    "   manejar la desproporcion entre BENIGN y tipos de ataque",
    "",
    "3. XGBoost como modelo principal: Con 92.8% de accuracy",
    "   sin optimizacion, tiene potencial para >99% con tuning",
    "",
    "4. Deploy en red real: Probar con trafico de ETH en vez de loopback",
    "   para capturar patrones reales de red",
    "",
    "5. Monitoreo continuo: Implementar re-entrenamiento periodico",
    "   para adaptarse a nuevos tipos de ataques",
])

# ==========================================
# REFERENCIAS
# ==========================================
add_section_slide("", "Referencias Bibliograficas")

add_content_slide("Referencias (APA 7)", [
    "[1] Becerra-Suarez, F. L., Fernandez-Roman, I., & Forero, M. G. (2024).",
    "    Improvement of DDoS Attack Detection through ML and Data Processing.",
    "    Mathematics, 12(9), 1294. https://doi.org/10.3390/math12091294",
    "",
    "[2] Sharafaldin, I., Lashkari, A. H., Hakak, S., & Ghorbani, A. A. (2019).",
    "    Developing Realistic DDoS Attack Dataset and Taxonomy. ACM ICCST.",
    "    https://doi.org/10.1145/3340997.3341005",
    "",
    "[3] Najam, N. R. & Abduljawad, R. A. (2023). RF-RFE-SMOTE: A DoS and",
    "    DDoS Attack Detection Framework. NTU-JET, 2(2), 29-47.",
    "    https://doi.org/10.56286/ntujet.v2i2.436",
    "",
    "[4] Ma, R., Chen, X., & Zhai, R. (2023). A DDoS Attack Detection Method",
    "    Based on Natural Selection of Features and Models. Electronics, 12(4),",
    "    1059. https://doi.org/10.3390/electronics12041059",
], os.path.join(BASE_DIR, 'resultados', 'matrices_confusion.png'))

add_content_slide("Referencias (APA 7) - Continuacion", [
    "[5] Ali, T. E., Chong, Y.-W., & Manickam, S. (2023). Machine Learning",
    "    Techniques to Detect a DDoS Attack in SDN: A Systematic Review.",
    "    Applied Sciences, 13(5), 3183. https://doi.org/10.3390/app13053183",
    "",
    "[6] Mualfah, D., Ardiansyah, R., & Gunawan, R. (2025). Classification of",
    "    DDoS attacks using random forest and class weight on CICDDoS2019.",
    "    Jurnal CoSciTech, 6(3), 530-535.",
    "    https://doi.org/10.37859/coscitech.v6i3.10731",
    "",
    "[7] Mittal, M., Kumar, K., & Behal, S. (2022). Deep learning approaches",
    "    for detecting DDoS attacks: a systematic review. Soft Computing.",
    "    https://doi.org/10.1007/s00500-021-06608-1",
    "",
    "[8] Cloudflare. (2024). DDoS threat report. Recuperado de",
    "    https://www.cloudflare.com/ddos/threat-report/",
])

# ==========================================
# SLIDE FINAL
# ==========================================
add_title_slide(
    "Gracias",
    "Christopher Abrego\nUniversidad Tecnologica de Panama\nIA Aplicada a la Ciberseguridad\n\nGitHub: github.com/Chrisavt/examen-ia-ddos-utp"
)

# ==========================================
# GUARDAR
# ==========================================
output_path = os.path.join(BASE_DIR, 'presentacion_ddos_utp.pptx')
prs.save(output_path)
print(f"Presentacion guardada: {output_path}")
print(f"Total de diapositivas: {len(prs.slides)}")
